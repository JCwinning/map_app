#!/usr/bin/env python3
"""
Create an enhanced PowerPoint presentation for the Shop Map Manager project.
Features modern design, gradients, visual hierarchy, and actual app screenshots.
Updated to include Supabase integration, cloud sync, and image upload features.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from os.path import join, dirname, abspath

# Color scheme - Modern gradient palette
COLORS = {
    'primary': RGBColor(41, 98, 255),      # Vibrant blue
    'secondary': RGBColor(99, 102, 241),    # Indigo
    'accent': RGBColor(236, 72, 153),       # Pink
    'success': RGBColor(34, 197, 94),       # Green
    'warning': RGBColor(251, 146, 60),      # Orange
    'dark': RGBColor(15, 23, 42),           # Dark slate
    'light': RGBColor(248, 250, 252),       # Light gray
    'white': RGBColor(255, 255, 255),
    'gradient_start': RGBColor(59, 130, 246),
    'gradient_end': RGBColor(147, 51, 234),
}

def add_gradient_background(slide, width=Inches(10), height=Inches(5.625)):
    """Add a subtle gradient background to the slide."""
    # Create a rectangle covering the whole slide
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        width, height
    )
    shape.fill.gradient()
    shape.fill.gradient_angle = 135
    shape.fill.gradient_stops[0].color.rgb = COLORS['light']
    shape.fill.gradient_stops[1].color.rgb = RGBColor(230, 235, 245)
    shape.line.fill.background()

    # Send to back
    shape.z_order = 0

def add_title_slide(prs, title, subtitle):
    """Add a stunning title slide with gradient background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add gradient background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.gradient()
    background.fill.gradient_angle = 135
    background.fill.gradient_stops[0].color.rgb = COLORS['gradient_start']
    background.fill.gradient_stops[1].color.rgb = COLORS['gradient_end']
    background.line.fill.background()

    # Add decorative circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7), Inches(-1),
        Inches(4), Inches(4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.font.color.brightness = 0.1
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_section_header_slide(prs, title, subtitle=""):
    """Add a section header slide with modern styling."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(1.5), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['primary']
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = COLORS['primary']

def add_content_slide(prs, title, content_items, icon=""):
    """Add a content slide with modern styling."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Subtle background
    add_gradient_background(slide)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.4),
        Inches(9), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['primary']
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.6), Inches(0.8))
    title_frame = title_box.text_frame
    title_text = f"{icon} {title}" if icon else title
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(3.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        p.font.color.rgb = COLORS['dark']

def add_features_slide(prs):
    """Add a features slide with visual cards."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_gradient_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 Key Features"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # Feature cards - reduced to 6 for better fit (2x3 grid)
    features = [
        ("👤 Authentication", "Sign up/Login with email via Supabase", COLORS['primary']),
        ("☁️ Cloud Sync", "Sync data across devices with Supabase", COLORS['secondary']),
        ("📍 Interactive Map", "Display shops on Gaode map tiles", COLORS['primary']),
        ("🔍 Smart Search", "Search via Gaode Map API", COLORS['secondary']),
        ("📸 Image Upload", "Upload & manage shop photos", COLORS['success']),
        ("✅ Visit Tracking", "Track visited vs want-to-visit", COLORS['accent']),
    ]

    box_width = Inches(3.0)
    box_height = Inches(1.1)
    start_left = Inches(0.5)
    start_top = Inches(1.8)
    h_gap = Inches(0.25)
    v_gap = Inches(0.25)

    for i, (feature_name, description, color) in enumerate(features):
        row = i // 3
        col = i % 3

        box_left = start_left + col * (box_width + h_gap)
        box_top = start_top + row * (box_height + v_gap)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left, box_top,
            box_width, box_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Feature icon/name
        name_box = slide.shapes.add_textbox(
            box_left + Inches(0.1), box_top + Inches(0.08),
            box_width - Inches(0.2), Inches(0.35)
        )
        name_frame = name_box.text_frame
        name_frame.text = feature_name
        name_para = name_frame.paragraphs[0]
        name_para.font.size = Pt(15)
        name_para.font.bold = True
        name_para.font.color.rgb = color

        # Description
        desc_box = slide.shapes.add_textbox(
            box_left + Inches(0.1), box_top + Inches(0.45),
            box_width - Inches(0.2), Inches(0.55)
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = description
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = COLORS['dark']

def add_screenshot_slide(prs, title, image_path, caption=""):
    """Add a slide with a screenshot and caption."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    add_gradient_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # Screenshot - with proper height constraint
    max_width = Inches(9)
    max_height = Inches(3.5)  # Limit height to prevent overflow

    # Add picture and get the shape
    pic = slide.shapes.add_picture(
        image_path,
        Inches(0.5), Inches(1),
        width=max_width
    )

    # Check if height exceeds max, if so scale down
    if pic.height > max_height:
        # Calculate scale factor
        scale_factor = max_height / pic.height
        new_width = pic.width * scale_factor
        # Re-add with correct dimensions
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            image_path,
            Inches(0.5), Inches(1),
            width=new_width
        )

    # Caption - position based on actual image height
    caption_top = Inches(1) + pic.height + Inches(0.15)
    if caption:
        caption_box = slide.shapes.add_textbox(Inches(0.5), caption_top, Inches(9), Inches(0.5))
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        caption_para = caption_frame.paragraphs[0]
        caption_para.font.size = Pt(16)
        caption_para.font.color.rgb = COLORS['primary']
        caption_para.alignment = PP_ALIGN.CENTER

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column slide with modern styling."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_gradient_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.4),
        Inches(4.2), Inches(0.5)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = COLORS['primary']
    left_header.line.fill.background()

    left_header_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(3.8), Inches(0.4))
    left_header_frame = left_header_box.text_frame
    left_header_frame.text = left_title
    left_header_para = left_header_frame.paragraphs[0]
    left_header_para.font.size = Pt(20)
    left_header_para.font.bold = True
    left_header_para.font.color.rgb = COLORS['white']

    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(4), Inches(3))
    left_content_frame = left_content.text_frame
    left_content_frame.word_wrap = True

    for i, item in enumerate(left_items):
        p = left_content_frame.add_paragraph() if i > 0 else left_content_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        p.font.color.rgb = COLORS['dark']

    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3), Inches(1.4),
        Inches(4.2), Inches(0.5)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = COLORS['secondary']
    right_header.line.fill.background()

    right_header_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(3.8), Inches(0.4))
    right_header_frame = right_header_box.text_frame
    right_header_frame.text = right_title
    right_header_para = right_header_frame.paragraphs[0]
    right_header_para.font.size = Pt(20)
    right_header_para.font.bold = True
    right_header_para.font.color.rgb = COLORS['white']

    # Right column content
    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(2.1), Inches(4), Inches(3))
    right_content_frame = right_content.text_frame
    right_content_frame.word_wrap = True

    for i, item in enumerate(right_items):
        p = right_content_frame.add_paragraph() if i > 0 else right_content_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        p.font.color.rgb = COLORS['dark']

def add_benefit_slide(prs):
    """Add a 'Why This Matters' benefit slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_gradient_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "💡 Why Use Shop Map Manager?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    title_para.alignment = PP_ALIGN.CENTER

    benefits = [
        ("Never Forget", "Keep track of interesting places you want to visit"),
        ("Personal Journal", "Build your personal map of experiences"),
        ("Smart Discovery", "Find new spots through intelligent search"),
        ("Visual Planning", "Plan visits with interactive map visualization"),
    ]

    box_width = Inches(4.2)
    box_height = Inches(1.2)
    start_left = Inches(0.55)
    start_top = Inches(1.6)
    gap = Inches(0.25)

    for i, (benefit_title, description) in enumerate(benefits):
        row = i // 2
        col = i % 2

        box_left = start_left + col * (box_width + gap)
        box_top = start_top + row * (box_height + gap)

        # Benefit card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            box_left, box_top,
            box_width, box_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['white']
        card.line.color.rgb = COLORS['primary']
        card.line.width = Pt(3)

        # Title
        title_b = slide.shapes.add_textbox(
            box_left + Inches(0.15), box_top + Inches(0.15),
            box_width - Inches(0.3), Inches(0.4)
        )
        title_b_frame = title_b.text_frame
        title_b_frame.text = benefit_title
        title_b_para = title_b_frame.paragraphs[0]
        title_b_para.font.size = Pt(18)
        title_b_para.font.bold = True
        title_b_para.font.color.rgb = COLORS['primary']

        # Description
        desc_b = slide.shapes.add_textbox(
            box_left + Inches(0.15), box_top + Inches(0.55),
            box_width - Inches(0.3), Inches(0.5)
        )
        desc_b_frame = desc_b.text_frame
        desc_b_frame.word_wrap = True
        desc_b_frame.text = description
        desc_b_para = desc_b_frame.paragraphs[0]
        desc_b_para.font.size = Pt(14)
        desc_b_para.font.color.rgb = COLORS['dark']

def create_presentation():
    """Create the complete enhanced presentation."""
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Get screenshot paths
    screenshot_dir = join(dirname(abspath(__file__)), 'static', 'screenshots')
    table_screenshot = join(screenshot_dir, 'table_view.png')
    map_screenshot = join(screenshot_dir, 'map_view.png')

    # Slide 1: Title
    add_title_slide(
        prs,
        "Shop Map Manager",
        "Your Personal Store Discovery & Location Tracking System"
    )

    # Slide 2: What is it?
    add_content_slide(prs, "What is Shop Map Manager?", [
        "🗺️ A powerful web app for discovering and tracking shops/locations",
        "🐍 Built with Python and Streamlit for easy deployment",
        "🇨🇳 Uses Gaode (高德) Maps for China location services",
        "☁️ Dual mode: Local (CSV) or Cloud (Supabase) storage",
        "👤 Optional auth for cross-device sync & image uploads",
        "🎯 Perfect for travelers, foodies, and urban explorers"
    ])

    # Slide 3: Why This Matters
    add_benefit_slide(prs)

    # Slide 4: Key Features
    add_features_slide(prs)

    # Slide 5: Screenshot - Table View
    add_screenshot_slide(
        prs,
        "📋 Data Table View",
        table_screenshot,
        "Intuitive spreadsheet-style editor with real-time data validation"
    )

    # Slide 6: Screenshot - Map View
    add_screenshot_slide(
        prs,
        "🗺️ Interactive Map View",
        map_screenshot,
        "Color-coded markers with detailed popups for each location"
    )

    # Slide 7: How It Works
    add_content_slide(prs, "🔄 How It Works", [
        "🔐 Choose mode: Local (no signup) or Cloud (login required)",
        "🔍 Search shops by name via Gaode API sidebar",
        "➕ Click 'Add to List' to save search results instantly",
        "✏️ Edit shop details in the data table view",
        "📸 Upload photos to shops (cloud mode only)",
        "🗺️ View & track locations on interactive map"
    ])

    # Slide 8: Technology Stack
    add_two_column_slide(
        prs,
        "⚙️ Technology Stack",
        "Frontend & UI",
        [
            "Python 3.x - Core language",
            "Streamlit - Modern web framework",
            "Folium - Interactive map visualization",
            "Streamlit-Folium - Seamless integration",
            "Custom HTML/CSS - Popup styling"
        ],
        "Backend & Data",
        [
            "Pandas - Powerful data manipulation",
            "Supabase - Authentication & Database",
            "Supabase Storage - Image hosting",
            "CSV - Local data storage option",
            "Gaode Map API - Location search",
            "Requests - HTTP API client"
        ]
    )

    # Slide 9: Architecture & Design
    add_content_slide(prs, "🏗️ Architecture & Design", [
        "📁 Modular: app.py (main) + map_utils.py (API & storage)",
        "💾 Dual storage: Local CSV or Supabase cloud",
        "🔐 Optional auth via Supabase (Email/Password)",
        "☁️ Auto cloud sync when logged in",
        "📸 Image storage via Supabase Storage",
        "🗺️ Gaode tiles for Chinese maps",
        "🎨 Color-coded: Red (Visited) / Green (Want to Visit)"
    ])

    # Slide 10: Data Structure
    add_two_column_slide(
        prs,
        "📊 Data Structure",
        "Core Information",
        [
            "User ID (cloud mode)",
            "Shop Name (required)",
            "City location",
            "Full address",
            "Latitude & Longitude",
            "Shop Type/Category"
        ],
        "User Tracking Fields",
        [
            "Journey Type (Coffee/Scenery/etc)",
            "Visit Status dropdown",
            "Rating (1-5 stars)",
            "Notes (free-form text)",
            "Image URLs (JSON array)"
        ]
    )

    # Slide 11: Getting Started
    add_section_header_slide(prs, "🚀 Getting Started")

    # Slide 12: Installation Steps
    add_content_slide(prs, "Installation & Setup", [
        "📦 Install: pip install -r requirements.txt",
        "🔑 Create .env file with API keys (see .env.example)",
        "🔑 Required: GAODE_API_KEY, SUPABASE_URL, SUPABASE_KEY",
        "▶️ Run: streamlit run app.py",
        "🌐 Open browser to: http://localhost:8501",
        "✨ Start searching for places!"
    ])

    # Slide 13: Future Enhancements
    add_content_slide(prs, "🔮 Future Enhancements", [
        "📱 Enhanced mobile-responsive design for on-the-go access",
        "🔔 Smart visit reminders and notifications",
        "🗺️ Route planning between multiple shops",
        "📊 Statistics dashboard with insights",
        "👥 Social sharing of favorite places",
        "🤖 AI-powered recommendations based on preferences"
    ])

    # Slide 14: Thank You
    add_title_slide(
        prs,
        "Thank You! 🎉",
        "Shop Map Manager - Your Personal Location Journal"
    )

    return prs

def main():
    """Main function to create and save the presentation."""
    print("Creating enhanced PowerPoint presentation...")
    prs = create_presentation()
    output_file = "Shop_Map_Manager_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Enhanced presentation saved to: {output_file}")
    print(f"✓ Features modern design, gradients, and actual screenshots")
    return output_file

if __name__ == "__main__":
    main()
